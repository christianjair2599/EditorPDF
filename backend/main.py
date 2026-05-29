from fastapi import FastAPI, File, UploadFile, Form, Request, Header, Depends, HTTPException, BackgroundTasks
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse, JSONResponse
from pydantic import BaseModel
import shutil, os, uuid, json, csv, re, time, hmac, hashlib
import stripe
import datetime
from database import SessionLocal, engine
from models import Base, Subscription

Base.metadata.create_all(bind=engine)

stripe.api_key = os.getenv("STRIPE_SECRET_KEY", "")
STRIPE_WEBHOOK_SECRET = os.getenv("STRIPE_WEBHOOK_SECRET", "")
STRIPE_PRICE_ID = os.getenv("STRIPE_PRICE_ID", "")
FRONTEND_URL = os.getenv("FRONTEND_URL", "https://editorpdf-christian-mayangas-projects.vercel.app")

# Lemon Squeezy
LEMON_API_KEY      = os.getenv("LEMONSQUEEZY_API_KEY", "")
LEMON_STORE_ID     = os.getenv("LEMONSQUEEZY_STORE_ID", "")
LEMON_VARIANT_ID   = os.getenv("LEMONSQUEEZY_VARIANT_ID", "")
LEMON_WEBHOOK_SECRET = os.getenv("LEMONSQUEEZY_WEBHOOK_SECRET", "")
from pdf2docx import Converter
import fitz  # PyMuPDF
from PIL import Image as PILImage
from docx import Document as DocxDocument
from docx.shared import Inches
from pptx import Presentation
from pptx.util import Inches as PptxInches, Pt
import openpyxl
from reportlab.lib.pagesizes import letter
from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
from reportlab.lib.styles import getSampleStyleSheet
from reportlab.lib import colors

app = FastAPI()

ALLOWED_ORIGINS = [
    origin.strip()
    for origin in os.getenv("ALLOWED_ORIGINS", "").split(",")
    if origin.strip()
] or [
    "https://editorpdf-christian-mayangas-projects.vercel.app",
    "https://docuflow-app.web.app",
    "http://localhost:3000",
    "http://localhost:3001",
]

app.add_middleware(
    CORSMiddleware,
    allow_origins=ALLOWED_ORIGINS,
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

UPLOAD_DIR = "uploads"
os.makedirs(UPLOAD_DIR, exist_ok=True)

# ── File cleanup ──────────────────────────────────────────────────────────────

_last_cleanup: float = 0

def cleanup_old_files():
    """Delete files older than 24 hours. Runs at most once per hour."""
    global _last_cleanup
    now = time.time()
    if now - _last_cleanup < 3600:
        return
    _last_cleanup = now
    cutoff = now - 86400  # 24 hours
    try:
        for filename in os.listdir(UPLOAD_DIR):
            filepath = os.path.join(UPLOAD_DIR, filename)
            if os.path.isfile(filepath) and os.path.getmtime(filepath) < cutoff:
                os.remove(filepath)
    except Exception as e:
        print(f"Cleanup error: {e}")


# ── Helpers ───────────────────────────────────────────────────────────────────

def hex_to_rgb(hex_color: str) -> tuple:
    hex_color = hex_color.lstrip("#")
    return (
        int(hex_color[0:2], 16) / 255,
        int(hex_color[2:4], 16) / 255,
        int(hex_color[4:6], 16) / 255,
    )

def safe_text(t: str) -> str:
    return t.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;")

def text_to_pdf(text: str, output_path: str):
    doc = SimpleDocTemplate(output_path, pagesize=letter,
                            leftMargin=50, rightMargin=50, topMargin=50, bottomMargin=50)
    styles = getSampleStyleSheet()
    style = styles["Normal"]
    style.fontSize = 11
    style.leading = 16
    story = []
    for line in text.split("\n"):
        line = line.strip()
        if line:
            try:
                story.append(Paragraph(safe_text(line), style))
            except Exception:
                pass
        else:
            story.append(Spacer(1, 8))
    if not story:
        story.append(Paragraph("(Sin contenido)", style))
    doc.build(story)

def table_to_pdf(rows: list[list], output_path: str):
    """Renders a list-of-rows table into a nicely styled PDF."""
    doc = SimpleDocTemplate(output_path, pagesize=letter,
                            leftMargin=40, rightMargin=40, topMargin=40, bottomMargin=40)
    styles = getSampleStyleSheet()
    if not rows:
        doc.build([Paragraph("(Sin contenido)", styles["Normal"])])
        return
    # Convert all cells to string
    data = [[str(c) if c is not None else "" for c in row] for row in rows]
    t = Table(data, repeatRows=1)
    t.setStyle(TableStyle([
        ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#1e3a5f")),
        ("TEXTCOLOR",  (0, 0), (-1, 0), colors.white),
        ("FONTNAME",   (0, 0), (-1, 0), "Helvetica-Bold"),
        ("FONTSIZE",   (0, 0), (-1, -1), 9),
        ("GRID",       (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
        ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f5f7fa")]),
        ("VALIGN",     (0, 0), (-1, -1), "MIDDLE"),
        ("PADDING",    (0, 0), (-1, -1), 5),
    ]))
    doc.build([t])

def docx_to_pdf_high_fidelity(input_path: str, output_path: str):
    """Converts a DOCX file to a PDF file with high fidelity styling and structured layout."""
    import re
    from docx import Document as DocxDocument
    from docx.table import Table as DocxTable
    from docx.text.paragraph import Paragraph as DocxParagraph
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle, PageBreak
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY

    def map_font_name(docx_font_name):
        if not docx_font_name:
            return None
        name = docx_font_name.lower()
        if "times" in name or "georgia" in name:
            return "Times-Roman"
        elif "courier" in name or "mono" in name:
            return "Courier"
        else:
            return "Helvetica"

    def rgb_to_hex(rgb_color):
        if rgb_color is None:
            return None
        try:
            return f"#{rgb_color[0]:02x}{rgb_color[1]:02x}{rgb_color[2]:02x}"
        except Exception:
            return None

    doc = SimpleDocTemplate(
        output_path, 
        pagesize=letter,
        leftMargin=40, 
        rightMargin=40, 
        topMargin=40, 
        bottomMargin=40
    )
    
    styles = getSampleStyleSheet()
    styles_cache = {"Normal": styles["Normal"]}
    
    story = []
    docx_doc = DocxDocument(input_path)
    
    for child in docx_doc.element.body.iterchildren():
        tag = child.tag
        
        # 1. Convert paragraphs
        if tag.endswith('p'):
            p = DocxParagraph(child, docx_doc)
            
            p_markup = ""
            for run in p.runs:
                r_text = run.text
                if not r_text:
                    continue
                r_escaped = r_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")
                
                font = run.font
                font_attrs = []
                
                if font.color and font.color.rgb:
                    c_hex = rgb_to_hex(font.color.rgb)
                    if c_hex:
                        font_attrs.append(f'color="{c_hex}"')
                if font.size:
                    try:
                        font_attrs.append(f'size="{font.size.pt}"')
                    except Exception:
                        pass
                if font.name:
                    f_name = map_font_name(font.name)
                    if f_name:
                        font_attrs.append(f'name="{f_name}"')
                        
                run_markup = r_escaped
                if font_attrs:
                    run_markup = f"<font {' '.join(font_attrs)}>{run_markup}</font>"
                    
                if run.bold:
                    run_markup = f"<b>{run_markup}</b>"
                if run.italic:
                    run_markup = f"<i>{run_markup}</i>"
                if run.underline:
                    run_markup = f"<u>{run_markup}</u>"
                    
                p_markup += run_markup
                
            # Fallback if runs are empty but text is not
            if not p_markup and p.text:
                p_markup = p.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")
                
            style_name = p.style.name if p.style else "Normal"
            
            # Prepend bullet to list items if missing
            is_bullet = "bullet" in style_name.lower() or "list" in style_name.lower()
            if is_bullet and not p_markup.strip().startswith(("&bull;", "•", "-", "*")):
                p_markup = "&bull;&nbsp;&nbsp;" + p_markup
                
            # Alignment mapping
            align = p.alignment
            if align is None:
                try:
                    align = p.style.paragraph_format.alignment
                except Exception:
                    align = None
                    
            rl_align = TA_LEFT
            if align == WD_ALIGN_PARAGRAPH.CENTER:
                rl_align = TA_CENTER
            elif align == WD_ALIGN_PARAGRAPH.RIGHT:
                rl_align = TA_RIGHT
            elif align == WD_ALIGN_PARAGRAPH.JUSTIFY:
                rl_align = TA_JUSTIFY
                
            # Spacing and Indents
            p_format = p.paragraph_format
            left_indent = 0
            if p_format.left_indent is not None:
                try:
                    left_indent = p_format.left_indent.pt
                except Exception:
                    pass
            space_before = 0
            if p_format.space_before is not None:
                try:
                    space_before = p_format.space_before.pt
                except Exception:
                    pass
            space_after = 6
            if p_format.space_after is not None:
                try:
                    space_after = p_format.space_after.pt
                except Exception:
                    pass
                    
            if not p_markup.strip():
                # Represent empty paragraph as visual space
                story.append(Spacer(1, 10))
                continue
                
            style_key = f"{style_name}_{rl_align}_{left_indent}_{space_before}_{space_after}"
            if style_key not in styles_cache:
                font_size = 11
                leading = 14
                bold_font = False
                font_name = "Helvetica"
                
                s_name = style_name.lower()
                if "title" in s_name:
                    font_size = 24
                    leading = 28
                    bold_font = True
                elif "subtitle" in s_name:
                    font_size = 16
                    leading = 20
                elif "heading 1" in s_name or "heading1" in s_name:
                    font_size = 18
                    leading = 22
                    bold_font = True
                    space_before = 12 if space_before == 0 else space_before
                    space_after = 6 if space_after == 6 else space_after
                elif "heading 2" in s_name or "heading2" in s_name:
                    font_size = 14
                    leading = 18
                    bold_font = True
                    space_before = 10 if space_before == 0 else space_before
                    space_after = 4 if space_after == 6 else space_after
                elif "heading 3" in s_name or "heading3" in s_name:
                    font_size = 12
                    leading = 15
                    bold_font = True
                    space_before = 8 if space_before == 0 else space_before
                    space_after = 4 if space_after == 6 else space_after
                elif "list" in s_name or "bullet" in s_name:
                    left_indent = left_indent or 20
                    space_after = 3 if space_after == 6 else space_after
                    
                r_style = ParagraphStyle(
                    name=f"Style_{style_name}_{rl_align}_{left_indent}_{len(styles_cache)}",
                    parent=styles_cache["Normal"],
                    fontName="Helvetica-Bold" if bold_font else font_name,
                    fontSize=font_size,
                    leading=leading,
                    alignment=rl_align,
                    leftIndent=left_indent,
                    spaceBefore=space_before,
                    spaceAfter=space_after
                )
                styles_cache[style_key] = r_style
            else:
                r_style = styles_cache[style_key]
                
            story.append(Paragraph(p_markup, r_style))
            
        # 2. Convert tables
        elif tag.endswith('tbl'):
            table = DocxTable(child, docx_doc)
            num_rows = len(table.rows)
            if num_rows == 0:
                continue
            num_cols = len(table.rows[0].cells)
            if num_cols == 0:
                continue
                
            # Compute column widths in points (usable page width = 532)
            col_widths = []
            first_row = table.rows[0]
            for cell in first_row.cells:
                w = cell.width
                if w and w > 0:
                    col_widths.append(w / 12700.0)
                else:
                    col_widths.append(0.0)
            
            usable_w = 532.0
            if all(w == 0 for w in col_widths):
                col_widths = [usable_w / num_cols] * num_cols
            else:
                non_zero = [w for w in col_widths if w > 0]
                avg_w = sum(non_zero) / len(non_zero) if non_zero else usable_w / num_cols
                col_widths = [w if w > 0 else avg_w for w in col_widths]
                
                total_w = sum(col_widths)
                if total_w > usable_w or total_w < 200.0:
                    scale = usable_w / total_w
                    col_widths = [w * scale for w in col_widths]
                    
            # Parse cell paragraphs
            grid_data = []
            for row in table.rows:
                row_data = []
                for cell in row.cells:
                    cell_flowables = []
                    for cp in cell.paragraphs:
                        cp_markup = ""
                        for crun in cp.runs:
                            cr_text = crun.text
                            if not cr_text:
                                continue
                            cr_escaped = cr_text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")
                            
                            cfont = crun.font
                            cfont_attrs = []
                            
                            if cfont.color and cfont.color.rgb:
                                cc_hex = rgb_to_hex(cfont.color.rgb)
                                if cc_hex:
                                    cfont_attrs.append(f'color="{cc_hex}"')
                            if cfont.size:
                                try:
                                    cfont_attrs.append(f'size="{cfont.size.pt}"')
                                except Exception:
                                    pass
                            if cfont.name:
                                cf_name = map_font_name(cfont.name)
                                if cf_name:
                                    cfont_attrs.append(f'name="{cf_name}"')
                                    
                            crun_markup = cr_escaped
                            if cfont_attrs:
                                crun_markup = f"<font {' '.join(cfont_attrs)}>{crun_markup}</font>"
                                
                            if crun.bold:
                                crun_markup = f"<b>{crun_markup}</b>"
                            if crun.italic:
                                crun_markup = f"<i>{crun_markup}</i>"
                            if crun.underline:
                                crun_markup = f"<u>{crun_markup}</u>"
                                
                            cp_markup += crun_markup
                            
                        if not cp_markup and cp.text:
                            cp_markup = cp.text.replace("&", "&amp;").replace("<", "&lt;").replace(">", "&gt;").replace("\n", "<br/>")
                            
                        if not cp_markup.strip():
                            continue
                            
                        cell_p_style = styles_cache.get("TableCellStyle")
                        if not cell_p_style:
                            cell_p_style = ParagraphStyle(
                                name="TableCellStyle",
                                parent=styles_cache["Normal"],
                                fontSize=9,
                                leading=11,
                                spaceAfter=2
                            )
                            styles_cache["TableCellStyle"] = cell_p_style
                            
                        cell_flowables.append(Paragraph(cp_markup, cell_p_style))
                        
                    if not cell_flowables:
                        cell_p_style = styles_cache.get("TableCellStyle")
                        if not cell_p_style:
                            cell_p_style = ParagraphStyle(
                                name="TableCellStyle",
                                parent=styles_cache["Normal"],
                                fontSize=9,
                                leading=11,
                                spaceAfter=2
                            )
                            styles_cache["TableCellStyle"] = cell_p_style
                        cell_flowables.append(Paragraph("", cell_p_style))
                        
                    row_data.append(cell_flowables)
                grid_data.append(row_data)
                
            ts = TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("PADDING", (0, 0), (-1, -1), 6),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eaeef3")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f7f9fb")]),
            ])
            
            rl_table = Table(grid_data, colWidths=col_widths)
            rl_table.setStyle(ts)
            story.append(rl_table)
            story.append(Spacer(1, 12))
            
    if not story:
        story.append(Paragraph("(Sin contenido)", styles["Normal"]))
        
    doc.build(story)


def stitch_images(imgs: list, output_path: str, fmt: str):
    if len(imgs) == 1:
        imgs[0].save(output_path, fmt)
        return
    total_w = max(i.width for i in imgs)
    total_h = sum(i.height for i in imgs)
    combined = PILImage.new("RGB", (total_w, total_h), (255, 255, 255))
    y = 0
    for img in imgs:
        combined.paste(img, (0, y))
        y += img.height
    combined.save(output_path, fmt)

def rows_to_json(rows: list[list]) -> str:
    if not rows:
        return "[]"
    headers = [str(h) if h else f"col{i}" for i, h in enumerate(rows[0])]
    result = []
    for row in rows[1:]:
        obj = {}
        for i, h in enumerate(headers):
            obj[h] = row[i] if i < len(row) else None
        result.append(obj)
    return json.dumps(result, ensure_ascii=False, indent=2)

def rows_to_xml(rows: list[list]) -> str:
    if not rows:
        return "<data/>"
    headers = [re.sub(r"[^a-zA-Z0-9_]", "_", str(h)) if h else f"col{i}"
               for i, h in enumerate(rows[0])]
    lines = ['<?xml version="1.0" encoding="UTF-8"?>', "<data>"]
    for row in rows[1:]:
        lines.append("  <row>")
        for i, h in enumerate(headers):
            val = row[i] if i < len(row) else ""
            lines.append(f"    <{h}>{safe_text(str(val))}</{h}>")
        lines.append("  </row>")
    lines.append("</data>")
    return "\n".join(lines)

def html_table(rows: list[list]) -> str:
    if not rows:
        return "<table></table>"
    header_cells = "".join(f"<th>{safe_text(str(c))}</th>" for c in rows[0])
    body_rows = ""
    for row in rows[1:]:
        cells = "".join(f"<td>{safe_text(str(c if c is not None else ''))}</td>" for c in row)
        body_rows += f"<tr>{cells}</tr>"
    return (
        "<table border='1' cellpadding='4' cellspacing='0' "
        f"style='border-collapse:collapse'><thead><tr>{header_cells}</tr></thead>"
        f"<tbody>{body_rows}</tbody></table>"
    )

def pdf_to_image_list(input_file: str) -> list:
    doc = fitz.open(input_file)
    imgs = []
    for page in doc:
        pix = page.get_pixmap(matrix=fitz.Matrix(2, 2))
        imgs.append(PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples))
    doc.close()
    return imgs

async def verify_premium_user(
    x_user_email: str = Header(None, alias="X-User-Email")
):
    db = SessionLocal()
    try:
        if not x_user_email:
            raise HTTPException(status_code=402, detail="Acceso Premium requerido. Por favor, inicia sesión.")
        
        # Check tester status first (QA / simulator)
        from models import User
        user = db.query(User).filter(User.email == x_user_email).first()
        if user and user.is_tester:
            return x_user_email
            
        # Check active subscription
        from models import Subscription
        sub = db.query(Subscription).filter(Subscription.email == x_user_email).first()
        is_premium = sub and sub.status == "active" and (
            sub.current_period_end is None or sub.current_period_end > datetime.datetime.utcnow()
        )
        if not is_premium:
            raise HTTPException(status_code=402, detail="Suscripción Premium requerida.")
    finally:
        db.close()
    return x_user_email

def remove_file(path: str):
    try:
        if os.path.exists(path):
            os.remove(path)
    except Exception as e:
        print(f"Error removing temp file {path}: {e}")

@app.get("/download/{filename}")
async def download_file(filename: str, background_tasks: BackgroundTasks, name: str = ""):
    safe = os.path.basename(filename)
    path = os.path.join(UPLOAD_DIR, safe)
    if not os.path.exists(path):
        return {"error": "Archivo no encontrado"}
    download_name = os.path.basename(name) if name else safe
    background_tasks.add_task(remove_file, path)
    return FileResponse(path, filename=download_name)


@app.post("/upload/")

async def upload_pdf(file: UploadFile = File(...)):
    cleanup_old_files()
    file_id = str(uuid.uuid4())
    path = os.path.join(UPLOAD_DIR, f"{file_id}_{file.filename}")
    with open(path, "wb") as buf:
        shutil.copyfileobj(file.file, buf)
    return {"message": "Archivo subido con éxito", "filename": file.filename, "file_path": path}


@app.post("/convert-any/")
async def convert_any(
    file: UploadFile = File(...),
    output_format: str = Form(...),
    x_user_email: str = Header(None, alias="X-User-Email")
):
    cleanup_old_files()
    file_id = str(uuid.uuid4())
    original_name = os.path.splitext(file.filename or "archivo")[0]
    original_ext = os.path.splitext(file.filename or "")[1].lower().lstrip(".")
    safe_name = f"{file_id}.{original_ext}" if original_ext else f"{file_id}.bin"
    input_file = os.path.join(UPLOAD_DIR, safe_name)

    with open(input_file, "wb") as buf:
        shutil.copyfileobj(file.file, buf)

    # Validar tamaño del archivo en backend
    file_size_mb = os.path.getsize(input_file) / (1024 * 1024)
    db = SessionLocal()
    is_premium = False
    try:
        if x_user_email:
            from models import User, Subscription
            user = db.query(User).filter(User.email == x_user_email).first()
            is_tester = user and user.is_tester
            sub = db.query(Subscription).filter(Subscription.email == x_user_email).first()
            is_premium = is_tester or (sub and sub.status == "active" and (
                sub.current_period_end is None or sub.current_period_end > datetime.datetime.utcnow()
            ))
    finally:
        db.close()

    max_allowed = 50 if is_premium else 5
    if file_size_mb > max_allowed:
        try:
            os.remove(input_file)
        except:
            pass
        return JSONResponse({"error": f"Límite de tamaño excedido ({max_allowed} MB). El plan gratuito tiene un límite de 5 MB."}, status_code=402)

    base = os.path.splitext(input_file)[0]

    fmt = output_format.lower()
    src = original_ext
    output_file = None

    try:
        # ── PDF → * ───────────────────────────────────────────────────────────
        if src == "pdf":
            if fmt == "docx":
                output_file = f"{base}.docx"
                cv = Converter(input_file)
                cv.convert(output_file)
                cv.close()

            elif fmt in ("jpg", "png", "webp", "bmp", "tiff"):
                output_file = f"{base}.{fmt}"
                imgs = pdf_to_image_list(input_file)
                pil_fmt = {"jpg": "JPEG", "png": "PNG", "webp": "WEBP", "bmp": "BMP", "tiff": "TIFF"}[fmt]
                stitch_images(imgs, output_file, pil_fmt)

            elif fmt == "txt":
                output_file = f"{base}.txt"
                doc = fitz.open(input_file)
                text = "\n".join(page.get_text() for page in doc)
                doc.close()
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(text)

            elif fmt == "html":
                output_file = f"{base}.html"
                doc = fitz.open(input_file)
                html = "<html><body>" + "".join(page.get_text("html") for page in doc) + "</body></html>"
                doc.close()
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(html)

            elif fmt in ("xlsx", "csv"):
                # Extract text lines into a spreadsheet / csv
                doc = fitz.open(input_file)
                rows = [["Página", "Texto"]]
                for page_num, page in enumerate(doc, 1):
                    for line in page.get_text().splitlines():
                        if line.strip():
                            rows.append([page_num, line.strip()])
                doc.close()
                if fmt == "csv":
                    output_file = f"{base}.csv"
                    with open(output_file, "w", newline="", encoding="utf-8") as f:
                        csv.writer(f).writerows(rows)
                else:
                    output_file = f"{base}.xlsx"
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    for row in rows:
                        ws.append(row)
                    wb.save(output_file)

            elif fmt == "json":
                doc = fitz.open(input_file)
                pages_data = []
                for page_num, page in enumerate(doc, 1):
                    pages_data.append({"page": page_num, "text": page.get_text()})
                doc.close()
                output_file = f"{base}.json"
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(pages_data, f, ensure_ascii=False, indent=2)

            elif fmt == "xml":
                doc = fitz.open(input_file)
                lines = ['<?xml version="1.0" encoding="UTF-8"?>', "<document>"]
                for page_num, page in enumerate(doc, 1):
                    lines.append(f'  <page number="{page_num}">')
                    for line in page.get_text().splitlines():
                        if line.strip():
                            lines.append(f"    <line>{safe_text(line.strip())}</line>")
                    lines.append("  </page>")
                lines.append("</document>")
                doc.close()
                output_file = f"{base}.xml"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write("\n".join(lines))

            elif fmt == "pptx":
                doc = fitz.open(input_file)
                prs = Presentation()
                slide_layout = prs.slide_layouts[1]
                for page_num, page in enumerate(doc):
                    slide = prs.slides.add_slide(slide_layout)
                    slide.shapes.title.text = f"Página {page_num + 1}"
                    tf = slide.placeholders[1].text_frame
                    tf.text = page.get_text()[:800]
                doc.close()
                output_file = f"{base}.pptx"
                prs.save(output_file)

        # ── DOCX / DOC → * ───────────────────────────────────────────────────
        elif src in ("docx", "doc"):
            docx = DocxDocument(input_file)
            paragraphs = [p.text for p in docx.paragraphs]
            full_text = "\n".join(paragraphs)

            if fmt == "pdf":
                output_file = f"{base}.pdf"
                docx_to_pdf_high_fidelity(input_file, output_file)

            elif fmt == "txt":
                output_file = f"{base}.txt"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(full_text)

            elif fmt == "html":
                output_file = f"{base}.html"
                body = "".join(f"<p>{safe_text(p)}</p>" for p in paragraphs if p.strip())
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"<html><body>{body}</body></html>")

            elif fmt in ("xlsx", "csv"):
                rows = [["#", "Párrafo"]] + [[i + 1, p] for i, p in enumerate(paragraphs) if p.strip()]
                if fmt == "csv":
                    output_file = f"{base}.csv"
                    with open(output_file, "w", newline="", encoding="utf-8") as f:
                        csv.writer(f).writerows(rows)
                else:
                    output_file = f"{base}.xlsx"
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    for row in rows:
                        ws.append(row)
                    wb.save(output_file)

            elif fmt == "json":
                output_file = f"{base}.json"
                data = [{"index": i + 1, "text": p} for i, p in enumerate(paragraphs) if p.strip()]
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)

            elif fmt == "xml":
                output_file = f"{base}.xml"
                items = "".join(
                    f"  <paragraph index='{i+1}'>{safe_text(p)}</paragraph>"
                    for i, p in enumerate(paragraphs) if p.strip()
                )
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f'<?xml version="1.0" encoding="UTF-8"?>\n<document>\n{items}\n</document>')

            elif fmt == "pptx":
                prs = Presentation()
                slide_layout = prs.slide_layouts[1]
                for p in paragraphs:
                    if p.strip():
                        slide = prs.slides.add_slide(slide_layout)
                        slide.shapes.title.text = p[:80]
                        prs.slides[-1].placeholders[1].text_frame.text = ""
                output_file = f"{base}.pptx"
                prs.save(output_file)

            elif fmt == "md":
                output_file = f"{base}.md"
                md_lines = []
                for p in paragraphs:
                    if p.strip():
                        md_lines.append(p.strip())
                        md_lines.append("")
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write("\n".join(md_lines))

        # ── PPTX / PPT → * ───────────────────────────────────────────────────
        elif src in ("pptx", "ppt"):
            prs = Presentation(input_file)
            slides_data = []
            for i, slide in enumerate(prs.slides, 1):
                texts = [shape.text.strip() for shape in slide.shapes
                         if hasattr(shape, "text") and shape.text.strip()]
                slides_data.append({"slide": i, "texts": texts})

            if fmt == "pdf":
                lines = []
                for s in slides_data:
                    lines.append(f"─── Diapositiva {s['slide']} ───")
                    lines.extend(s["texts"])
                    lines.append("")
                output_file = f"{base}.pdf"
                text_to_pdf("\n".join(lines), output_file)

            elif fmt == "txt":
                output_file = f"{base}.txt"
                lines = []
                for s in slides_data:
                    lines.append(f"=== Diapositiva {s['slide']} ===")
                    lines.extend(s["texts"])
                    lines.append("")
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write("\n".join(lines))

            elif fmt == "html":
                output_file = f"{base}.html"
                sections = ""
                for s in slides_data:
                    paras = "".join(f"<p>{safe_text(t)}</p>" for t in s["texts"])
                    sections += f"<section><h2>Diapositiva {s['slide']}</h2>{paras}</section>"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"<html><body>{sections}</body></html>")

            elif fmt == "docx":
                output_file = f"{base}.docx"
                doc = DocxDocument()
                for s in slides_data:
                    doc.add_heading(f"Diapositiva {s['slide']}", level=1)
                    for t in s["texts"]:
                        doc.add_paragraph(t)
                doc.save(output_file)

            elif fmt in ("xlsx", "csv"):
                rows = [["Diapositiva", "Texto"]]
                for s in slides_data:
                    for t in s["texts"]:
                        rows.append([s["slide"], t])
                if fmt == "csv":
                    output_file = f"{base}.csv"
                    with open(output_file, "w", newline="", encoding="utf-8") as f:
                        csv.writer(f).writerows(rows)
                else:
                    output_file = f"{base}.xlsx"
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    for row in rows:
                        ws.append(row)
                    wb.save(output_file)

            elif fmt == "json":
                output_file = f"{base}.json"
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(slides_data, f, ensure_ascii=False, indent=2)

            elif fmt == "xml":
                output_file = f"{base}.xml"
                lines = ['<?xml version="1.0" encoding="UTF-8"?>', "<presentation>"]
                for s in slides_data:
                    lines.append(f'  <slide number="{s["slide"]}">')
                    for t in s["texts"]:
                        lines.append(f"    <text>{safe_text(t)}</text>")
                    lines.append("  </slide>")
                lines.append("</presentation>")
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write("\n".join(lines))

        # ── XLSX / XLS → * ───────────────────────────────────────────────────
        elif src in ("xlsx", "xls"):
            wb = openpyxl.load_workbook(input_file, data_only=True)
            ws = wb.active
            rows = list(ws.iter_rows(values_only=True))

            if fmt == "csv":
                output_file = f"{base}.csv"
                with open(output_file, "w", newline="", encoding="utf-8") as f:
                    csv.writer(f).writerows(
                        [str(c) if c is not None else "" for c in row] for row in rows
                    )

            elif fmt == "txt":
                output_file = f"{base}.txt"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write("\n".join(
                        "\t".join(str(c) if c is not None else "" for c in row) for row in rows
                    ))

            elif fmt == "pdf":
                output_file = f"{base}.pdf"
                table_to_pdf(rows, output_file)

            elif fmt == "html":
                output_file = f"{base}.html"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"<html><body>{html_table(rows)}</body></html>")

            elif fmt == "json":
                output_file = f"{base}.json"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(rows_to_json(rows))

            elif fmt == "xml":
                output_file = f"{base}.xml"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(rows_to_xml(rows))

            elif fmt == "docx":
                output_file = f"{base}.docx"
                doc = DocxDocument()
                table = doc.add_table(rows=len(rows), cols=len(rows[0]) if rows else 1)
                table.style = "Table Grid"
                for r_idx, row in enumerate(rows):
                    for c_idx, cell in enumerate(row):
                        table.rows[r_idx].cells[c_idx].text = str(cell) if cell is not None else ""
                doc.save(output_file)

            elif fmt == "pptx":
                prs = Presentation()
                layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(layout)
                rows_clean = [[str(c) if c is not None else "" for c in row] for row in rows[:20]]
                if rows_clean:
                    n_rows = len(rows_clean)
                    n_cols = max(len(r) for r in rows_clean)
                    tbl = slide.shapes.add_table(n_rows, n_cols,
                        PptxInches(0.5), PptxInches(1), PptxInches(9), PptxInches(5)).table
                    for r_idx, row in enumerate(rows_clean):
                        for c_idx in range(n_cols):
                            tbl.cell(r_idx, c_idx).text = row[c_idx] if c_idx < len(row) else ""
                output_file = f"{base}.pptx"
                prs.save(output_file)

        # ── CSV → * ──────────────────────────────────────────────────────────
        elif src == "csv":
            with open(input_file, "r", encoding="utf-8", errors="replace") as f:
                rows = list(csv.reader(f))

            if fmt == "xlsx":
                output_file = f"{base}.xlsx"
                wb = openpyxl.Workbook()
                ws = wb.active
                for row in rows:
                    ws.append(row)
                wb.save(output_file)

            elif fmt == "txt":
                output_file = f"{base}.txt"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write("\n".join("\t".join(row) for row in rows))

            elif fmt == "pdf":
                output_file = f"{base}.pdf"
                table_to_pdf(rows, output_file)

            elif fmt == "html":
                output_file = f"{base}.html"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"<html><body>{html_table(rows)}</body></html>")

            elif fmt == "json":
                output_file = f"{base}.json"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(rows_to_json(rows))

            elif fmt == "xml":
                output_file = f"{base}.xml"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(rows_to_xml(rows))

            elif fmt == "docx":
                output_file = f"{base}.docx"
                doc = DocxDocument()
                if rows:
                    table = doc.add_table(rows=len(rows), cols=len(rows[0]))
                    table.style = "Table Grid"
                    for r_idx, row in enumerate(rows):
                        for c_idx, cell in enumerate(row):
                            table.rows[r_idx].cells[c_idx].text = cell
                doc.save(output_file)

            elif fmt == "pptx":
                prs = Presentation()
                layout = prs.slide_layouts[5]
                slide = prs.slides.add_slide(layout)
                rows_slice = rows[:20]
                if rows_slice:
                    n_rows = len(rows_slice)
                    n_cols = max(len(r) for r in rows_slice)
                    tbl = slide.shapes.add_table(n_rows, n_cols,
                        PptxInches(0.5), PptxInches(1), PptxInches(9), PptxInches(5)).table
                    for r_idx, row in enumerate(rows_slice):
                        for c_idx in range(n_cols):
                            tbl.cell(r_idx, c_idx).text = row[c_idx] if c_idx < len(row) else ""
                output_file = f"{base}.pptx"
                prs.save(output_file)

        # ── Images → * ───────────────────────────────────────────────────────
        elif src in ("jpg", "jpeg", "png", "webp", "bmp", "gif", "tiff"):
            img_mode = "RGBA" if fmt == "png" else "RGB"
            img = PILImage.open(input_file).convert(img_mode)

            pil_fmts = {
                "pdf":  ("pdf",  "PDF"),
                "jpg":  ("jpg",  "JPEG"),
                "jpeg": ("jpg",  "JPEG"),
                "png":  ("png",  "PNG"),
                "webp": ("webp", "WEBP"),
                "bmp":  ("bmp",  "BMP"),
                "tiff": ("tiff", "TIFF"),
                "gif":  ("gif",  "GIF"),
                "ico":  ("ico",  "ICO"),
            }
            if fmt not in pil_fmts:
                return {"error": f"Conversión de imagen a '{fmt}' no soportada."}
            ext_out, pil_fmt = pil_fmts[fmt]
            output_file = f"{base}.{ext_out}"
            if fmt == "ico":
                img_rgb = img.convert("RGBA")
                img_rgb = img_rgb.resize((256, 256), PILImage.LANCZOS)
                img_rgb.save(output_file, "ICO", sizes=[(256, 256), (128, 128), (64, 64), (32, 32), (16, 16)])
            elif fmt == "pdf":
                img.convert("RGB").save(output_file, "PDF")
            else:
                save_kwargs = {}
                if fmt in ("jpg", "jpeg"):
                    save_kwargs = {"quality": 92}
                elif fmt == "webp":
                    save_kwargs = {"quality": 90}
                img.save(output_file, pil_fmt, **save_kwargs)

        # ── TXT → * ──────────────────────────────────────────────────────────
        elif src == "txt":
            with open(input_file, "r", encoding="utf-8", errors="replace") as f:
                text = f.read()
            lines = text.splitlines()

            if fmt == "pdf":
                output_file = f"{base}.pdf"
                text_to_pdf(text, output_file)

            elif fmt == "docx":
                output_file = f"{base}.docx"
                doc = DocxDocument()
                for line in lines:
                    doc.add_paragraph(line)
                doc.save(output_file)

            elif fmt == "html":
                output_file = f"{base}.html"
                body = "".join(f"<p>{safe_text(l)}</p>" for l in lines if l.strip())
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f"<html><body>{body}</body></html>")

            elif fmt in ("xlsx", "csv"):
                rows = [["#", "Línea"]] + [[i + 1, l] for i, l in enumerate(lines) if l.strip()]
                if fmt == "csv":
                    output_file = f"{base}.csv"
                    with open(output_file, "w", newline="", encoding="utf-8") as f:
                        csv.writer(f).writerows(rows)
                else:
                    output_file = f"{base}.xlsx"
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    for row in rows:
                        ws.append(row)
                    wb.save(output_file)

            elif fmt == "json":
                output_file = f"{base}.json"
                data = [{"line": i + 1, "text": l} for i, l in enumerate(lines) if l.strip()]
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)

            elif fmt == "xml":
                output_file = f"{base}.xml"
                items = "".join(
                    f'  <line number="{i+1}">{safe_text(l)}</line>'
                    for i, l in enumerate(lines) if l.strip()
                )
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f'<?xml version="1.0" encoding="UTF-8"?>\n<text>\n{items}\n</text>')

            elif fmt == "pptx":
                prs = Presentation()
                layout = prs.slide_layouts[1]
                for line in lines:
                    if line.strip():
                        slide = prs.slides.add_slide(layout)
                        slide.shapes.title.text = line[:80]
                        slide.placeholders[1].text_frame.text = ""
                output_file = f"{base}.pptx"
                prs.save(output_file)

            elif fmt == "md":
                output_file = f"{base}.md"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(text)

        # ── HTML / HTM → * ───────────────────────────────────────────────────
        elif src in ("html", "htm"):
            with open(input_file, "r", encoding="utf-8", errors="replace") as f:
                raw_html = f.read()
            plain = re.sub(r"<[^>]+>", " ", raw_html)
            plain = re.sub(r"\s+", " ", plain).strip()
            lines = [l.strip() for l in plain.split(".") if l.strip()]

            if fmt == "txt":
                output_file = f"{base}.txt"
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(plain)

            elif fmt == "pdf":
                output_file = f"{base}.pdf"
                text_to_pdf(plain, output_file)

            elif fmt == "docx":
                output_file = f"{base}.docx"
                doc = DocxDocument()
                for line in lines:
                    if line:
                        doc.add_paragraph(line + ".")
                doc.save(output_file)

            elif fmt == "json":
                output_file = f"{base}.json"
                data = {"content": plain, "sentences": lines}
                with open(output_file, "w", encoding="utf-8") as f:
                    json.dump(data, f, ensure_ascii=False, indent=2)

            elif fmt == "xml":
                output_file = f"{base}.xml"
                items = "".join(f"  <sentence>{safe_text(l)}.</sentence>" for l in lines if l)
                with open(output_file, "w", encoding="utf-8") as f:
                    f.write(f'<?xml version="1.0" encoding="UTF-8"?>\n<html_content>\n{items}\n</html_content>')

            elif fmt in ("xlsx", "csv"):
                rows = [["#", "Oración"]] + [[i + 1, l + "."] for i, l in enumerate(lines) if l]
                if fmt == "csv":
                    output_file = f"{base}.csv"
                    with open(output_file, "w", newline="", encoding="utf-8") as f:
                        csv.writer(f).writerows(rows)
                else:
                    output_file = f"{base}.xlsx"
                    wb = openpyxl.Workbook()
                    ws = wb.active
                    for row in rows:
                        ws.append(row)
                    wb.save(output_file)

        # ── Unsupported ────────────────────────────────────────────────────────
        if output_file is None:
            return {"error": f"Conversión de '{src}' a '{fmt}' no soportada."}

    except Exception as e:
        return {"error": f"Error en la conversión: {str(e)}"}

    friendly_name = f"{original_name}.{fmt}"
    return {"message": "Conversión exitosa", "output_file": os.path.basename(output_file), "friendly_name": friendly_name}


def url_to_pdf_high_fidelity(url: str, soup, output_path: str):
    """Converts a BeautifulSoup web page object to a highly styled PDF document."""
    from reportlab.lib.pagesizes import letter
    from reportlab.platypus import SimpleDocTemplate, Paragraph, Spacer, Table, TableStyle
    from reportlab.lib.styles import getSampleStyleSheet, ParagraphStyle
    from reportlab.lib import colors
    from reportlab.lib.enums import TA_LEFT, TA_CENTER, TA_RIGHT, TA_JUSTIFY

    doc = SimpleDocTemplate(
        output_path, 
        pagesize=letter,
        leftMargin=45, 
        rightMargin=45, 
        topMargin=45, 
        bottomMargin=45
    )
    
    styles = getSampleStyleSheet()
    styles_cache = {"Normal": styles["Normal"]}
    
    story = []
    
    # 1. Header (Title of Page and Source Link)
    title = soup.title.string.strip() if soup.title and soup.title.string else "Sitio Web"
    
    title_style = ParagraphStyle(
        name="URLTitle",
        parent=styles["Normal"],
        fontName="Helvetica-Bold",
        fontSize=20,
        leading=24,
        textColor=colors.HexColor("#1e3a5f"),
        spaceAfter=6
    )
    story.append(Paragraph(safe_text(title), title_style))
    
    source_style = ParagraphStyle(
        name="URLSource",
        parent=styles["Normal"],
        fontName="Helvetica-Oblique",
        fontSize=9,
        leading=12,
        textColor=colors.HexColor("#666666"),
        spaceAfter=15
    )
    story.append(Paragraph(f"Fuente: <a href='{url}'><font color='#0066cc'>{safe_text(url)}</font></a>", source_style))
    story.append(Spacer(1, 8))
    
    # Track which elements are already parsed (e.g. elements inside tables to avoid double parsing)
    in_table_elements = set()
    for table_el in soup.find_all("table"):
        for sub_el in table_el.find_all(["h1", "h2", "h3", "h4", "p", "li", "blockquote", "pre", "td", "th"]):
            in_table_elements.add(id(sub_el))
            
    # Traverse page elements
    for el in soup.find_all(["h1", "h2", "h3", "h4", "p", "li", "blockquote", "pre", "table"]):
        if id(el) in in_table_elements:
            continue
            
        tag = el.name
        
        # A. TABLE
        if tag == "table":
            rows_data = []
            for tr in el.find_all("tr"):
                row_cells = []
                for cell in tr.find_all(["td", "th"]):
                    cell_text = cell.get_text(strip=True)
                    cell_style = styles_cache.get("URLTableCell")
                    if not cell_style:
                        cell_style = ParagraphStyle(
                            name="URLTableCell",
                            parent=styles["Normal"],
                            fontSize=9,
                            leading=11,
                            spaceAfter=2
                        )
                        styles_cache["URLTableCell"] = cell_style
                    
                    if cell.name == "th":
                        cell_style = ParagraphStyle(
                            name=f"URLTableHead_{len(styles_cache)}",
                            parent=cell_style,
                            fontName="Helvetica-Bold"
                        )
                    row_cells.append(Paragraph(safe_text(cell_text), cell_style))
                if row_cells:
                    rows_data.append(row_cells)
                    
            if not rows_data:
                continue
                
            num_cols = max(len(r) for r in rows_data)
            col_w = 522.0 / num_cols
            col_widths = [col_w] * num_cols
            
            for r in rows_data:
                while len(r) < num_cols:
                    r.append(Paragraph("", styles_cache.get("URLTableCell")))
                    
            ts = TableStyle([
                ("GRID", (0, 0), (-1, -1), 0.5, colors.HexColor("#cccccc")),
                ("VALIGN", (0, 0), (-1, -1), "TOP"),
                ("PADDING", (0, 0), (-1, -1), 5),
                ("BACKGROUND", (0, 0), (-1, 0), colors.HexColor("#eaeef3")),
                ("ROWBACKGROUNDS", (0, 1), (-1, -1), [colors.white, colors.HexColor("#f7f9fb")]),
            ])
            rl_table = Table(rows_data, colWidths=col_widths)
            rl_table.setStyle(ts)
            story.append(rl_table)
            story.append(Spacer(1, 10))
            
        # B. TEXT TAGS
        else:
            text = el.get_text(strip=True)
            if not text:
                continue
                
            escaped = safe_text(text)
            
            font_size = 11
            leading = 14
            bold = False
            italic = False
            textColor = colors.HexColor("#333333")
            space_before = 0
            space_after = 6
            left_indent = 0
            
            if tag == "h1":
                font_size = 18
                leading = 22
                bold = True
                textColor = colors.HexColor("#1e3a5f")
                space_before = 12
                space_after = 6
            elif tag == "h2":
                font_size = 15
                leading = 18
                bold = True
                textColor = colors.HexColor("#2c3e50")
                space_before = 10
                space_after = 5
            elif tag == "h3":
                font_size = 13
                leading = 16
                bold = True
                space_before = 8
                space_after = 4
            elif tag == "h4":
                font_size = 11
                leading = 14
                bold = True
                space_before = 6
                space_after = 3
            elif tag == "blockquote":
                left_indent = 20
                italic = True
                textColor = colors.HexColor("#555555")
                space_after = 8
            elif tag == "pre":
                font_size = 9
                leading = 11
                textColor = colors.HexColor("#000000")
                space_after = 8
            elif tag == "li":
                left_indent = 15
                escaped = "&bull;&nbsp;&nbsp;" + escaped
                space_after = 3
                
            style_key = f"{tag}_{bold}_{italic}_{left_indent}"
            if style_key not in styles_cache:
                r_style = ParagraphStyle(
                    name=f"Style_{style_key}_{len(styles_cache)}",
                    parent=styles["Normal"],
                    fontName="Courier" if tag == "pre" else ("Helvetica-Bold" if bold else ("Helvetica-Oblique" if italic else "Helvetica")),
                    fontSize=font_size,
                    leading=leading,
                    textColor=textColor,
                    leftIndent=left_indent,
                    spaceBefore=space_before,
                    spaceAfter=space_after
                )
                styles_cache[style_key] = r_style
            else:
                r_style = styles_cache[style_key]
                
            story.append(Paragraph(escaped, r_style))
            
    if not story:
        story.append(Paragraph("(Sin contenido)", styles["Normal"]))
        
    doc.build(story)


# ── URL → PDF / Image / TXT / HTML ───────────────────────────────────────────
@app.post("/url-to-pdf/")
async def url_to_format(url: str = Form(...), output_format: str = Form("pdf")):
    import re as _re
    import requests as _req
    from bs4 import BeautifulSoup

    if not _re.match(r"^https?://", url):
        return JSONResponse({"error": "URL inválida. Debe comenzar con http:// o https://"}, status_code=400)

    file_id = str(uuid.uuid4())
    fmt = output_format.lower()
    headers = {"User-Agent": "Mozilla/5.0 (Windows NT 10.0; Win64; x64) AppleWebKit/537.36"}

    try:
        resp = _req.get(url, timeout=30, headers=headers)
        resp.raise_for_status()
        soup = BeautifulSoup(resp.text, "html.parser")
        for tag in soup(["script", "style", "nav", "footer", "iframe", "noscript"]):
            tag.decompose()

        title = soup.title.string.strip() if soup.title and soup.title.string else url
        paragraphs = [title, f"Fuente: {url}", ""]
        for el in soup.find_all(["h1","h2","h3","h4","p","li","td","th","blockquote","pre"]):
            t = el.get_text(separator=" ", strip=True)
            if t:
                paragraphs.append(t)

        plain_text = "\n".join(paragraphs)

        if fmt in ("pdf", "png", "jpg", "jpeg"):
            pdf_path = os.path.join(UPLOAD_DIR, f"{file_id}_tmp.pdf")
            try:
                from urllib.parse import quote
                microlink_url = f"https://api.microlink.io/?url={quote(url)}&pdf=true&embed=pdf.url"
                api_resp = _req.get(microlink_url, timeout=35)
                if api_resp.status_code == 200 and len(api_resp.content) > 1000:
                    with open(pdf_path, "wb") as f:
                        f.write(api_resp.content)
                else:
                    raise Exception(f"Microlink returned status {api_resp.status_code}")
            except Exception as e:
                print(f"Microlink print failed, falling back to local layout engine: {e}")
                url_to_pdf_high_fidelity(url, soup, pdf_path)

            if fmt == "pdf":
                out_path = pdf_path
            else:
                pil_fmt = {"png": "PNG", "jpg": "JPEG", "jpeg": "JPEG"}[fmt]
                out_path = os.path.join(UPLOAD_DIR, f"{file_id}.{fmt}")
                imgs = pdf_to_image_list(pdf_path)
                stitch_images(imgs, out_path, pil_fmt)
                os.remove(pdf_path)

        elif fmt == "txt":
            out_path = os.path.join(UPLOAD_DIR, f"{file_id}.txt")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(plain_text)

        elif fmt == "html":
            out_path = os.path.join(UPLOAD_DIR, f"{file_id}.html")
            with open(out_path, "w", encoding="utf-8") as f:
                f.write(resp.text)

        else:
            return JSONResponse({"error": f"Formato '{fmt}' no soportado para URLs"}, status_code=400)

    except _req.exceptions.ConnectionError:
        return JSONResponse({"error": "No se pudo conectar a la URL. Verifica que el enlace sea accesible."}, status_code=400)
    except _req.exceptions.Timeout:
        return JSONResponse({"error": "La URL tardó demasiado en responder."}, status_code=400)
    except _req.exceptions.HTTPError as e:
        return JSONResponse({"error": f"La URL devolvió un error HTTP: {e.response.status_code}"}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": f"Error al procesar la URL: {str(e)}"}, status_code=500)

    # Build a friendly filename from the URL hostname
    from urllib.parse import urlparse
    host = urlparse(url).hostname or "pagina"
    host = host.replace("www.", "")
    friendly_name = f"{host}.{fmt}"
    return {"message": "Conversión exitosa", "output_file": os.path.basename(out_path), "friendly_name": friendly_name}


# ── Legacy /convert/ ─────────────────────────────────────────────────────────
@app.post("/convert/")
async def convert_pdf(file: UploadFile = File(...), output_format: str = Form("docx")):
    return await convert_any(file=file, output_format=output_format)


# ── PDF Editor endpoints ──────────────────────────────────────────────────────

@app.post("/edit/")
async def edit_pdf(
    file: UploadFile = File(...),
    comment: str = Form("Editado con EditorPDF"),
    x: float = Form(72), y: float = Form(72),
    font_size: float = Form(12), color: str = Form("#FF0000"),
    apply_to_all: bool = Form(True), page_num: int = Form(0),
):
    file_id = str(uuid.uuid4())
    input_file = os.path.join(UPLOAD_DIR, f"{file_id}_{file.filename}")
    with open(input_file, "wb") as buf:
        shutil.copyfileobj(file.file, buf)
    try:
        rgb = hex_to_rgb(color)
        doc = fitz.open(input_file)
        pages = range(len(doc)) if apply_to_all else [page_num]
        for i in pages:
            if 0 <= i < len(doc):
                doc[i].insert_text((x, y), comment, fontsize=font_size, color=rgb)
        base, _ = os.path.splitext(input_file)
        output_file = f"{base}_edited.pdf"
        doc.save(output_file)
        doc.close()
    except Exception as e:
        return {"error": f"Error al editar PDF: {str(e)}"}
    return {"message": "Edición exitosa", "output_file": os.path.basename(output_file)}


@app.post("/extract/")
async def extract_text(file: UploadFile = File(...)):
    file_id = str(uuid.uuid4())
    input_file = os.path.join(UPLOAD_DIR, f"{file_id}_{file.filename}")
    with open(input_file, "wb") as buf:
        shutil.copyfileobj(file.file, buf)
    try:
        doc = fitz.open(input_file)
        blocks = []
        for page_num in range(len(doc)):
            page = doc[page_num]
            page_dict = page.get_text("dict")
            for block in page_dict.get("blocks", []):
                if block.get("type") != 0:
                    continue
                text_parts, font_size, color = [], 12, "#000000"
                for line in block.get("lines", []):
                    for span in line.get("spans", []):
                        t = span.get("text", "")
                        if t.strip():
                            text_parts.append(t)
                        font_size = round(span.get("size", 12))
                        c = span.get("color", 0)
                        if isinstance(c, int):
                            color = f"#{(c >> 16) & 0xFF:02x}{(c >> 8) & 0xFF:02x}{c & 0xFF:02x}"
                full_text = " ".join(text_parts).strip()
                if not full_text:
                    continue
                bb = block["bbox"]
                blocks.append({
                    "id": f"p{page_num}_b{block.get('number', len(blocks))}",
                    "page": page_num,
                    "x0": round(bb[0]), "y0": round(bb[1]),
                    "x1": round(bb[2]), "y1": round(bb[3]),
                    "text": full_text, "font_size": font_size, "color": color,
                })
        doc.close()
        return {"blocks": blocks, "file_id": file_id}
    except Exception as e:
        return {"error": f"Error al extraer texto: {str(e)}"}


@app.post("/edit-blocks/")
async def edit_blocks(file_id: str = Form(...), edits: str = Form(...)):
    files = [f for f in os.listdir(UPLOAD_DIR) if f.startswith(file_id)]
    if not files:
        return {"error": "Archivo no encontrado. Vuelve a cargar el PDF."}
    input_file = os.path.join(UPLOAD_DIR, files[0])
    try:
        edits_list = json.loads(edits)
        doc = fitz.open(input_file)
        for edit in edits_list:
            page = doc[edit["page"]]
            rect = fitz.Rect(edit["x0"], edit["y0"], edit["x1"], edit["y1"])
            shape = page.new_shape()
            shape.draw_rect(rect)
            shape.finish(fill=(1, 1, 1), color=(1, 1, 1), width=0)
            shape.commit()
            page.insert_textbox(
                rect, edit["new_text"],
                fontsize=float(edit.get("font_size", 12)),
                fontname=edit.get("font", "helv"),
                color=hex_to_rgb(edit.get("color", "#000000")),
                align=0,
            )
        base, _ = os.path.splitext(input_file)
        doc.save(output_file)
        doc.close()
    except Exception as e:
        return {"error": f"Error al aplicar ediciones: {str(e)}"}
    return {"message": "Ediciones aplicadas", "output_file": os.path.basename(output_file)}

@app.post("/improve-text/")
async def improve_text(
    text: str = Form(...),
    x_user_email: str = Depends(verify_premium_user)
):
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {"error": "ANTHROPIC_API_KEY no configurada en el servidor."}

    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key)
        response = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=1024,
            messages=[{"role": "user", "content": (
                "Mejora el siguiente texto de un documento, haciéndolo más claro, "
                "profesional y bien redactado. Mantén el mismo significado e idioma. "
                "Responde ÚNICAMENTE con el texto mejorado, sin explicaciones ni comillas:\n\n"
                + text
            )}],
        )
        return {"improved_text": response.content[0].text}
    except ImportError:
        return {"error": "Paquete 'anthropic' no instalado en el servidor."}
    except Exception as e:
        return {"error": f"Error con IA: {str(e)}"}


# ── PDF Utilities (Merge, Split, Compress) ────────────────────────────────

@app.post("/merge-pdf/")
async def merge_pdf(
    files: list[UploadFile] = File(...),
    x_user_email: str = Header(None, alias="X-User-Email")
):
    """Merge multiple PDF files into one."""
    if len(files) < 2:
        return {"error": "Se necesitan al menos 2 PDFs para fusionar."}

    if len(files) > 2:
        db = SessionLocal()
        is_premium = False
        try:
            if x_user_email:
                from models import User, Subscription
                user = db.query(User).filter(User.email == x_user_email).first()
                is_tester = user and user.is_tester
                sub = db.query(Subscription).filter(Subscription.email == x_user_email).first()
                is_premium = is_tester or (sub and sub.status == "active" and (
                    sub.current_period_end is None or sub.current_period_end > datetime.datetime.utcnow()
                ))
        finally:
            db.close()
        if not is_premium:
            return JSONResponse({"error": "El plan gratuito solo admite fusionar hasta 2 PDFs. Actualiza a Premium para fusionar ilimitados."}, status_code=402)


    try:
        temp_paths = []

        # Save all files first
        for file in files:
            if not file.filename.lower().endswith(".pdf"):
                return {"error": f"'{file.filename}' no es un PDF."}

            temp_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{file.filename}")
            content = await file.read()
            with open(temp_path, "wb") as f:
                f.write(content)
            temp_paths.append(temp_path)

        # Now merge using fitz instead of PyPDF2
        output_name = f"{uuid.uuid4()}_merged.pdf"
        output_path = os.path.join(UPLOAD_DIR, output_name)

        merged_doc = fitz.open()

        for temp_path in temp_paths:
            try:
                doc = fitz.open(temp_path)
                merged_doc.insert_pdf(doc)
                doc.close()
            except Exception as e:
                print(f"Error merging {temp_path}: {e}")
                raise

        merged_doc.save(output_path)
        merged_doc.close()

        # Cleanup
        for path in temp_paths:
            try:
                os.remove(path)
            except:
                pass

        return {"message": "PDFs fusionados", "output_file": output_name}
    except Exception as e:
        print(f"Merge error: {str(e)}")
        return {"error": f"Error al fusionar: {str(e)}"}


@app.post("/split-pdf/")
async def split_pdf(file: UploadFile = File(...), start_page: int = Form(0), end_page: int = Form(-1)):
    """Extract pages from a PDF (0-indexed)."""
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Solo se aceptan archivos PDF."}

    try:
        temp_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{file.filename}")
        content = await file.read()
        with open(temp_path, "wb") as f:
            f.write(content)

        # Use PyMuPDF for splitting
        doc = fitz.open(temp_path)
        total_pages = len(doc)

        if start_page < 0:
            start_page = 0
        if end_page < 0 or end_page >= total_pages:
            end_page = total_pages - 1

        if start_page > end_page:
            doc.close()
            os.remove(temp_path)
            return {"error": f"Páginas inválidas: {start_page}-{end_page}"}

        # Create new PDF with selected pages
        new_doc = fitz.open()
        for page_num in range(start_page, min(end_page + 1, total_pages)):
            new_doc.insert_pdf(doc, from_page=page_num, to_page=page_num)

        base = os.path.splitext(file.filename)[0]
        output_name = f"{uuid.uuid4()}_{base}_pages_{start_page+1}-{end_page+1}.pdf"
        output_path = os.path.join(UPLOAD_DIR, output_name)

        new_doc.save(output_path)
        new_doc.close()
        doc.close()
        os.remove(temp_path)

        return {"message": "PDF dividido", "output_file": output_name}
    except Exception as e:
        print(f"Split error: {str(e)}")
        return {"error": f"Error al dividir: {str(e)}"}


@app.post("/compress-pdf/")
async def compress_pdf(file: UploadFile = File(...), quality: int = Form(2)):
    """Compress PDF or image files."""
    if not file.filename:
        return {"error": "Nombre de archivo vacío."}

    try:
        temp_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{file.filename}")
        content = await file.read()
        with open(temp_path, "wb") as f:
            f.write(content)

        ext = os.path.splitext(file.filename)[1].lower().lstrip(".")
        base = os.path.splitext(file.filename)[0]
        output_name = f"{uuid.uuid4()}_{base}_compressed.{ext}"
        output_path = os.path.join(UPLOAD_DIR, output_name)

        # PDF compression
        if ext == "pdf":
            doc = fitz.open(temp_path)
            # Use garbage collection and deflate compression
            doc.save(output_path, garbage=4, deflate=True)
            doc.close()

        # Image compression
        elif ext in ("jpg", "jpeg", "png", "webp", "bmp"):
            img = PILImage.open(temp_path)

            # Quality levels: 0=no reduction, 1=10%, 2=20%, 3=30%
            scale = max(0.5, 1 - (quality * 0.1))
            new_width = max(50, int(img.width * scale))
            new_height = max(50, int(img.height * scale))
            img_resized = img.resize((new_width, new_height), PILImage.LANCZOS)

            if ext == "png":
                img_resized.save(output_path, "PNG", optimize=True)
            else:
                quality_jpeg = [95, 85, 75, 60][min(quality, 3)]
                img_resized.save(output_path, "JPEG", quality=quality_jpeg, optimize=True)

        else:
            # For other files, just copy
            shutil.copy(temp_path, output_path)

        os.remove(temp_path)
        return {"message": "Archivo comprimido", "output_file": output_name}
    except Exception as e:
        print(f"Compress error: {str(e)}")
        return {"error": f"Error al comprimir: {str(e)}"}


@app.post("/ocr-pdf/")
async def ocr_pdf(
    file: UploadFile = File(...),
    x_user_email: str = Depends(verify_premium_user)
):
    """Extract text from images or scanned PDFs using a serverless-friendly hybrid OCR system (local text extract + OCR.space API fallback)."""
    import requests
    temp_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{file.filename}")
    try:
        content = await file.read()
        with open(temp_path, "wb") as f:
            f.write(content)

        all_text = []
        fname = (file.filename or "").lower()

        if fname.endswith(".pdf"):
            doc = fitz.open(temp_path)
            for page_idx, page in enumerate(doc):
                page_num = page_idx + 1
                # 1. Intentar extracción de texto nativa del PDF
                native_text = page.get_text().strip()
                if len(native_text) > 30:
                    all_text.append(f"--- PÁGINA {page_num} (Texto Extraído) ---\n{native_text}")
                    continue

                # 2. Si no hay texto nativo, renderizar la página y enviarla a OCR.space
                try:
                    pix = page.get_pixmap(matrix=fitz.Matrix(2.0, 2.0))
                    page_img_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_page_{page_num}.jpg")
                    img = PILImage.frombytes("RGB", [pix.width, pix.height], pix.samples)
                    img.save(page_img_path, "JPEG")

                    # Consultar OCR.space
                    api_key = os.getenv("OCR_SPACE_API_KEY", "helloworld")
                    url = "https://api.ocr.space/parse/image"
                    with open(page_img_path, "rb") as img_file:
                        files = {"file": (os.path.basename(page_img_path), img_file, "image/jpeg")}
                        payload = {
                            "apikey": api_key,
                            "language": "spa",
                            "isOverlayRequired": False,
                        }
                        res = requests.post(url, files=files, data=payload, timeout=30)
                        res.raise_for_status()
                        data = res.json()

                    if os.path.exists(page_img_path):
                        os.remove(page_img_path)

                    if data.get("IsErroredOnProcessing"):
                        err_msg = data.get("ErrorMessage") or "Error de OCR en esta página"
                        all_text.append(f"--- PÁGINA {page_num} ---\n[Error de OCR: {err_msg}]")
                    else:
                        parsed = data.get("ParsedResults", [])
                        text = parsed[0].get("ParsedText", "").strip() if parsed else ""
                        all_text.append(f"--- PÁGINA {page_num} (OCR Scanned) ---\n{text}")

                except Exception as page_err:
                    all_text.append(f"--- PÁGINA {page_num} ---\n[Error al procesar página: {str(page_err)}]")
            doc.close()
        else:
            # Procesar imagen directa con OCR.space
            try:
                api_key = os.getenv("OCR_SPACE_API_KEY", "helloworld")
                url = "https://api.ocr.space/parse/image"
                with open(temp_path, "rb") as img_file:
                    files = {"file": (os.path.basename(temp_path), img_file)}
                    payload = {
                        "apikey": api_key,
                        "language": "spa",
                        "isOverlayRequired": False,
                    }
                    res = requests.post(url, files=files, data=payload, timeout=30)
                    res.raise_for_status()
                    data = res.json()

                if data.get("IsErroredOnProcessing"):
                    err_msg = data.get("ErrorMessage") or "Error de OCR"
                    all_text.append(f"--- DOCUMENTO ---\n[Error de OCR: {err_msg}]")
                else:
                    parsed = data.get("ParsedResults", [])
                    text = parsed[0].get("ParsedText", "").strip() if parsed else ""
                    all_text.append(f"--- DOCUMENTO (OCR Extraído) ---\n{text}")
            except Exception as img_err:
                all_text.append(f"--- DOCUMENTO ---\n[Error en OCR: {str(img_err)}]")

        output_name = f"{uuid.uuid4()}_ocr.txt"
        output_path = os.path.join(UPLOAD_DIR, output_name)
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n\n".join(all_text))

        return {"message": "OCR completado", "output_file": output_name}

    except Exception as e:
        print(f"OCR error: {str(e)}")
        return {"error": f"Error en OCR general: {str(e)}"}
    finally:
        if os.path.exists(temp_path):
            os.remove(temp_path)


@app.post("/watermark-pdf/")
async def watermark_pdf(
    file: UploadFile = File(...),
    text: str = Form("CONFIDENCIAL"),
    opacity: float = Form(0.3),
    font_size: int = Form(48),
    color: str = Form("#FF0000"),
    x_user_email: str = Depends(verify_premium_user)
):

    """Add a diagonal text watermark to every page of a PDF."""
    if not file.filename.lower().endswith(".pdf"):
        return {"error": "Solo se aceptan archivos PDF."}

    try:
        import math

        temp_path = os.path.join(UPLOAD_DIR, f"{uuid.uuid4()}_{file.filename}")
        content = await file.read()
        with open(temp_path, "wb") as f:
            f.write(content)

        doc = fitz.open(temp_path)

        # Parse color
        r, g, b = hex_to_rgb(color.lstrip("#") if not color.startswith("#") else color)

        for page in doc:
            pw, ph = page.rect.width, page.rect.height
            # Place watermark text diagonally across the page
            center_x = pw / 2
            center_y = ph / 2
            angle = -45  # degrees

            page.insert_text(
                fitz.Point(center_x - font_size * len(text) * 0.3, center_y),
                text,
                fontsize=font_size,
                color=(r, g, b),
                rotate=angle,
                overlay=True,
            )

        base = os.path.splitext(file.filename)[0]
        output_name = f"{uuid.uuid4()}_{base}_watermark.pdf"
        output_path = os.path.join(UPLOAD_DIR, output_name)
        doc.save(output_path, garbage=4, deflate=True)
        doc.close()
        os.remove(temp_path)

        return {"message": "Marca de agua añadida", "output_file": output_name}
    except Exception as e:
        print(f"Watermark error: {str(e)}")
        return {"error": f"Error al añadir marca de agua: {str(e)}"}


@app.post("/share/")
async def share_file(
    file: UploadFile = File(...),
    x_user_email: str = Depends(verify_premium_user)
):

    """Store a file and return a shareable token."""
    try:
        token = str(uuid.uuid4())
        ext = os.path.splitext(file.filename)[1]
        stored_name = f"{token}{ext}"
        stored_path = os.path.join(UPLOAD_DIR, stored_name)
        content = await file.read()
        with open(stored_path, "wb") as f:
            f.write(content)
        return {"token": token, "filename": stored_name}
    except Exception as e:
        return {"error": f"Error al compartir: {str(e)}"}


# ── Stripe / Payments ─────────────────────────────────────────────────────────

def get_or_create_subscription(db, email: str) -> Subscription:
    sub = db.query(Subscription).filter(Subscription.email == email).first()
    if not sub:
        sub = Subscription(email=email, status="free")
        db.add(sub)
        db.commit()
        db.refresh(sub)
    return sub


@app.post("/create-checkout-session")
async def create_checkout_session(request: Request):
    try:
        body = await request.json()
        email = body.get("email")
        if not email:
            return JSONResponse({"error": "Email requerido"}, status_code=400)
        if not stripe.api_key:
            return JSONResponse({"error": "Stripe no configurado"}, status_code=500)

        db = SessionLocal()
        sub = get_or_create_subscription(db, email)

        # Reuse existing Stripe customer if we have one
        customer_id = sub.stripe_customer_id
        if not customer_id:
            customer = stripe.Customer.create(email=email)
            customer_id = customer.id
            sub.stripe_customer_id = customer_id
            db.commit()
        db.close()

        session = stripe.checkout.Session.create(
            customer=customer_id,
            payment_method_types=["card"],
            line_items=[{"price": STRIPE_PRICE_ID, "quantity": 1}],
            mode="subscription",
            success_url=f"{FRONTEND_URL}/pricing?success=true&session_id={{CHECKOUT_SESSION_ID}}",
            cancel_url=f"{FRONTEND_URL}/pricing?canceled=true",
            metadata={"email": email},
        )
        return {"url": session.url}
    except stripe.StripeError as e:
        return JSONResponse({"error": str(e)}, status_code=400)
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.get("/subscription-status")
async def subscription_status(email: str):
    if not email:
        return {"status": "free", "isPremium": False, "isAdmin": False, "isTester": False}
    db = SessionLocal()
    try:
        from models import User
        user = db.query(User).filter(User.email == email).first()
        if not user:
            # Dynamically register the Google Auth user in the users table!
            user = User(email=email, usage_count=0)
            db.add(user)
            db.commit()
            db.refresh(user)
        is_admin = user.is_admin if user else False
        is_tester = user.is_tester if user else False

        # If it's a tester, they get premium unlocked automatically
        if user and user.is_tester:
            return {
                "status": "active",
                "isPremium": True,
                "currentPeriodEnd": None,
                "isAdmin": is_admin,
                "isTester": is_tester,
            }

        sub = get_or_create_subscription(db, email)
        is_premium = sub.status == "active" and (
            sub.current_period_end is None or sub.current_period_end > datetime.datetime.utcnow()
        )
        return {
            "status": sub.status,
            "isPremium": is_premium,
            "currentPeriodEnd": sub.current_period_end.isoformat() if sub.current_period_end else None,
            "isAdmin": is_admin,
            "isTester": is_tester,
        }
    finally:
        db.close()


@app.post("/customer-portal")
async def customer_portal(request: Request):
    try:
        body = await request.json()
        email = body.get("email")
        if not email:
            return JSONResponse({"error": "Email requerido"}, status_code=400)
        db = SessionLocal()
        sub = db.query(Subscription).filter(Subscription.email == email).first()
        db.close()
        if not sub or not sub.stripe_customer_id:
            return JSONResponse({"error": "Sin suscripción activa"}, status_code=404)
        portal = stripe.billing_portal.Session.create(
            customer=sub.stripe_customer_id,
            return_url=f"{FRONTEND_URL}/pricing",
        )
        return {"url": portal.url}
    except stripe.StripeError as e:
        return JSONResponse({"error": str(e)}, status_code=400)


@app.post("/stripe-webhook")
async def stripe_webhook(request: Request, stripe_signature: str = Header(None)):
    payload = await request.body()
    try:
        event = stripe.Webhook.construct_event(payload, stripe_signature, STRIPE_WEBHOOK_SECRET)
    except stripe.SignatureVerificationError:
        return JSONResponse({"error": "Invalid signature"}, status_code=400)

    db = SessionLocal()
    try:
        if event["type"] == "checkout.session.completed":
            session = event["data"]["object"]
            email = session.get("metadata", {}).get("email") or session.get("customer_details", {}).get("email")
            stripe_sub_id = session.get("subscription")
            customer_id = session.get("customer")
            if email and stripe_sub_id:
                sub = get_or_create_subscription(db, email)
                sub.stripe_customer_id = customer_id
                sub.stripe_subscription_id = stripe_sub_id
                sub.status = "active"
                db.commit()

        elif event["type"] in ("customer.subscription.updated", "customer.subscription.deleted"):
            stripe_sub = event["data"]["object"]
            customer_id = stripe_sub.get("customer")
            sub = db.query(Subscription).filter(Subscription.stripe_customer_id == customer_id).first()
            if sub:
                sub.status = stripe_sub.get("status")  # "active", "canceled", "past_due"
                period_end = stripe_sub.get("current_period_end")
                if period_end:
                    sub.current_period_end = datetime.datetime.utcfromtimestamp(period_end)
                sub.updated_at = datetime.datetime.utcnow()
                db.commit()
    finally:
        db.close()

    return {"received": True}


# ── Lemon Squeezy ─────────────────────────────────────────────────────────────

@app.post("/lemon-checkout")
async def lemon_checkout(request: Request):
    """Return a Lemon Squeezy checkout URL with email pre-filled."""
    try:
        body = await request.json()
        email = body.get("email", "")
        variant_id = LEMON_VARIANT_ID
        if not variant_id:
            return JSONResponse({"error": "Lemon Squeezy no configurado"}, status_code=500)
        import urllib.parse
        base = f"https://docufloww.lemonsqueezy.com/checkout/buy/{variant_id}"
        params = urllib.parse.urlencode({
            "checkout[email]": email,
            "checkout[custom][user_email]": email,
        })
        url = f"{base}?{params}"
        return {"url": url}
    except Exception as e:
        return JSONResponse({"error": str(e)}, status_code=500)


@app.post("/lemon-webhook")
async def lemon_webhook(request: Request, x_signature: str = Header(None, alias="X-Signature")):
    """Handle Lemon Squeezy webhook events."""
    payload = await request.body()

    # Verify signature
    if LEMON_WEBHOOK_SECRET and LEMON_WEBHOOK_SECRET != "pending":
        expected = hmac.new(
            LEMON_WEBHOOK_SECRET.encode("utf-8"), payload, hashlib.sha256
        ).hexdigest()
        if not hmac.compare_digest(expected, x_signature or ""):
            return JSONResponse({"error": "Invalid signature"}, status_code=400)

    try:
        data = json.loads(payload)
    except Exception:
        return JSONResponse({"error": "Invalid JSON"}, status_code=400)

    event_name = data.get("meta", {}).get("event_name", "")
    attrs = data.get("data", {}).get("attributes", {})
    custom = data.get("meta", {}).get("custom_data", {})
    email = custom.get("user_email") or attrs.get("user_email") or attrs.get("customer_email", "")

    if not email:
        # Try to get from order user email
        email = attrs.get("user_email", "")

    db = SessionLocal()
    try:
        if event_name in ("order_created", "subscription_created") and email:
            sub = get_or_create_subscription(db, email)
            sub.status = "active"
            ends_at = attrs.get("ends_at") or attrs.get("renews_at")
            if ends_at:
                try:
                    sub.current_period_end = datetime.datetime.fromisoformat(ends_at.replace("Z", "+00:00")).replace(tzinfo=None)
                except Exception:
                    pass
            sub.updated_at = datetime.datetime.utcnow()
            db.commit()

        elif event_name == "subscription_updated" and email:
            sub = db.query(Subscription).filter(Subscription.email == email).first()
            if sub:
                status = attrs.get("status", "")
                sub.status = "active" if status == "active" else status
                ends_at = attrs.get("ends_at") or attrs.get("renews_at")
                if ends_at:
                    try:
                        sub.current_period_end = datetime.datetime.fromisoformat(ends_at.replace("Z", "+00:00")).replace(tzinfo=None)
                    except Exception:
                        pass
                sub.updated_at = datetime.datetime.utcnow()
                db.commit()

        elif event_name in ("subscription_cancelled", "subscription_expired") and email:
            sub = db.query(Subscription).filter(Subscription.email == email).first()
            if sub:
                sub.status = "canceled"
                sub.updated_at = datetime.datetime.utcnow()
                db.commit()
    finally:
        db.close()

    return {"received": True}


# ── Admin & QA Panel Endpoints ───────────────────────────────────────────────

async def verify_admin_user(
    x_user_email: str = Header(None, alias="X-User-Email")
):
    db = SessionLocal()
    try:
        if not x_user_email:
            raise HTTPException(status_code=403, detail="Acceso denegado. Inicia sesión.")
        from models import User
        user = db.query(User).filter(User.email == x_user_email).first()
        if not user or not user.is_admin:
            raise HTTPException(status_code=403, detail="Acceso denegado. Requiere privilegios de administrador.")
    finally:
        db.close()
    return x_user_email

@app.get("/admin/stats")
async def admin_stats(x_user_email: str = Depends(verify_admin_user)):
    db = SessionLocal()
    try:
        from models import User, Subscription
        from sqlalchemy import func

        total_users = db.query(User).count()
        premium_users = db.query(Subscription).filter(Subscription.status == "active").count()
        
        # System usage stats: sum of usage_count across all users
        total_ops = db.query(func.sum(User.usage_count)).scalar() or 0
        
        # Calculate conversion rate
        conversion_rate = round((premium_users / total_users * 100), 1) if total_users > 0 else 0.0
        
        # Estimated MRR (Monthly Recurring Revenue) - premium users * $9 USD
        mrr = premium_users * 9

        return {
            "total_users": total_users,
            "premium_users": premium_users,
            "total_operations": total_ops,
            "conversion_rate": conversion_rate,
            "mrr": mrr,
        }
    finally:
        db.close()

@app.get("/admin/users")
async def admin_users(
    q: str = "",
    x_user_email: str = Depends(verify_admin_user)
):
    db = SessionLocal()
    try:
        from models import User, Subscription
        
        # Query users, support search filter
        query = db.query(User)
        if q:
            query = query.filter(User.email.contains(q))
        users = query.order_by(User.id.desc()).limit(100).all()
        
        # Map user records with their subscription details
        results = []
        for u in users:
            sub = db.query(Subscription).filter(Subscription.email == u.email).first()
            results.append({
                "id": u.id,
                "email": u.email,
                "usage_count": u.usage_count,
                "is_admin": u.is_admin,
                "is_tester": u.is_tester,
                "premium_status": sub.status if sub else "free",
                "current_period_end": sub.current_period_end.isoformat() if sub and sub.current_period_end else None
            })
        return results
    finally:
        db.close()

@app.post("/admin/users/{user_id}/toggle-role")
async def toggle_role(
    user_id: int,
    role: str = Form(...),  # "is_admin" | "is_tester"
    x_user_email: str = Depends(verify_admin_user)
):
    db = SessionLocal()
    try:
        from models import User
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            raise HTTPException(status_code=404, detail="Usuario no encontrado")
            
        # Protect against self-demotion
        if user.email == x_user_email and role == "is_admin":
            raise HTTPException(status_code=400, detail="No puedes quitarte el rol de administrador a ti mismo.")
            
        if role == "is_admin":
            user.is_admin = not user.is_admin
        elif role == "is_tester":
            user.is_tester = not user.is_tester
        else:
            raise HTTPException(status_code=400, detail="Rol inválido")
            
        db.commit()
        return {"success": True, "is_admin": user.is_admin, "is_tester": user.is_tester}
    finally:
        db.close()

@app.post("/admin/users/{user_id}/set-premium")
async def set_premium(
    user_id: int,
    status: str = Form(...),  # "active" | "free"
    x_user_email: str = Depends(verify_admin_user)
):
    db = SessionLocal()
    try:
        from models import User, Subscription
        user = db.query(User).filter(User.id == user_id).first()
        if not user:
            raise HTTPException(status_code=404, detail="Usuario no encontrado")
            
        sub = db.query(Subscription).filter(Subscription.email == user.email).first()
        if not sub:
            sub = Subscription(email=user.email, status="free")
            db.add(sub)
            db.commit()
            db.refresh(sub)
            
        sub.status = status
        if status == "active":
            # Set period end to 1 year in the future for convenience
            sub.current_period_end = datetime.datetime.utcnow() + datetime.timedelta(days=365)
        else:
            sub.current_period_end = None
            
        sub.updated_at = datetime.datetime.utcnow()
        db.commit()
        return {"success": True, "premium_status": sub.status}
    finally:
        db.close()

# ── Tester Simulator Endpoints ───────────────────────────────────────────────

async def verify_tester_user(
    x_user_email: str = Header(None, alias="X-User-Email")
):
    db = SessionLocal()
    try:
        if not x_user_email:
            raise HTTPException(status_code=403, detail="Acceso denegado. Inicia sesión.")
        from models import User
        user = db.query(User).filter(User.email == x_user_email).first()
        if not user or not user.is_tester:
            raise HTTPException(status_code=403, detail="Acceso denegado. Requiere perfil de tester de pruebas.")
    finally:
        db.close()
    return x_user_email

@app.post("/test/toggle-premium")
async def test_toggle_premium(x_user_email: str = Depends(verify_tester_user)):
    db = SessionLocal()
    try:
        from models import Subscription
        sub = db.query(Subscription).filter(Subscription.email == x_user_email).first()
        if not sub:
            sub = Subscription(email=x_user_email, status="free")
            db.add(sub)
            db.commit()
            db.refresh(sub)
            
        new_status = "active" if sub.status != "active" else "free"
        sub.status = new_status
        if new_status == "active":
            sub.current_period_end = datetime.datetime.utcnow() + datetime.timedelta(days=30)
        else:
            sub.current_period_end = None
            
        sub.updated_at = datetime.datetime.utcnow()
        db.commit()
        return {"success": True, "premium_status": sub.status}
    finally:
        db.close()

@app.post("/test/reset-limits")
async def test_reset_limits(x_user_email: str = Depends(verify_tester_user)):
    db = SessionLocal()
    try:
        from models import User
        user = db.query(User).filter(User.email == x_user_email).first()
        if user:
            user.usage_count = 0
            db.commit()
            return {"success": True, "usage_count": 0}
        raise HTTPException(status_code=404, detail="Usuario no encontrado")
    finally:
        db.close()

@app.post("/test/mock-webhook")
async def test_mock_webhook(
    event_type: str = Form(...),  # "order_created" | "subscription_cancelled"
    x_user_email: str = Depends(verify_tester_user)
):
    """Trigger a local mock Lemon Squeezy webhook to test background DB update logic instantly."""
    db = SessionLocal()
    try:
        from models import Subscription
        sub = db.query(Subscription).filter(Subscription.email == x_user_email).first()
        if not sub:
            sub = Subscription(email=x_user_email, status="free")
            db.add(sub)
            db.commit()
            db.refresh(sub)

        if event_type == "order_created":
            sub.status = "active"
            sub.current_period_end = datetime.datetime.utcnow() + datetime.timedelta(days=30)
        elif event_type == "subscription_cancelled":
            sub.status = "canceled"
            
        sub.updated_at = datetime.datetime.utcnow()
        db.commit()
        return {"success": True, "premium_status": sub.status, "event_simulated": event_type}
    finally:
        db.close()


class SupportChatRequest(BaseModel):
    message: str
    history: list[dict] = []

@app.post("/support-chat/")
async def support_chat(req: SupportChatRequest):
    api_key = os.getenv("ANTHROPIC_API_KEY")
    if not api_key:
        return {"reply": "¡Hola! Soy DocuBot, tu asistente de soporte de DocuFlow. (Nota: ANTHROPIC_API_KEY no configurada en el servidor). ¿En qué puedo ayudarte hoy?"}

    try:
        import anthropic
        client = anthropic.Anthropic(api_key=api_key)
        
        system_prompt = (
            "Eres DocuBot, el asistente inteligente de soporte técnico y comercial de DocuFlow.\n"
            "DocuFlow es un software SaaS premium e inteligente para editar, convertir y optimizar documentos PDF.\n"
            "Debes ser amable, conciso, profesional y ayudar al usuario en español.\n"
            "Información Clave de DocuFlow:\n"
            "- Plan Gratuito: 4 operaciones al día, archivos hasta 5 MB de tamaño, herramientas básicas (conversión de un archivo, editar texto, compresión, fusión hasta 2 PDFs).\n"
            "- Plan Premium: Operaciones ilimitadas, archivos hasta 50 MB, OCR avanzado (para extraer texto de escaneos e imágenes), mejora de textos con IA (Claude), procesamiento por lote, y compartir con enlace público.\n"
            "- Precios en Perú: S/ 34 al mes (moneda local). En Europa: 8 € al mes. En USA / Resto del mundo: $9 USD al mes.\n"
            "- Privacidad y Seguridad: Todos los documentos subidos están seguros encriptados y se eliminan automáticamente del disco del servidor inmediatamente después de ser descargados.\n"
            "- Si el usuario experimenta fallos o límites diarios, aliéntalo con educación a activar la Suscripción Premium en modo demo (completamente gratis para probar durante 7 días en la página de Precios).\n"
            "- Mantén tus respuestas breves y directas (máximo 3 oraciones por mensaje para que quepa bien en el chat flotante)."
        )

        formatted_messages = []
        for msg in req.history:
            role = msg.get("role", "user")
            content = msg.get("content", "")
            if role in ["user", "assistant"] and content:
                formatted_messages.append({"role": role, "content": content})
                
        formatted_messages.append({"role": "user", "content": req.message})

        response = client.messages.create(
            model="claude-haiku-4-5-20251001",
            max_tokens=400,
            system=system_prompt,
            messages=formatted_messages
        )
        reply = response.content[0].text
        return {"reply": reply}
    except Exception as e:
        return {"reply": f"Lo siento, he tenido un inconveniente técnico al procesar tu solicitud: {str(e)}"}
